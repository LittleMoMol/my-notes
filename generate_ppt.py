#!/usr/bin/env python3
"""从 第 1 章 引言.md 生成 PowerPoint 演示文稿"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# ── 主题色 ──
PRIMARY = RGBColor(0x1A, 0x47, 0x8A)       # 深蓝 - 主色调
ACCENT = RGBColor(0x2E, 0x8B, 0x57)        # 绿色 - 强调
DARK = RGBColor(0x2C, 0x3E, 0x50)          # 深灰
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)         # 浅灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MATH_BG = RGBColor(0xF8, 0xF9, 0xFA)       # 公式浅底色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 辅助函数 ──

def add_bg(slide, color=LIGHT):
    """给幻灯片添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color=WHITE, left=0, top=0, width=None, height=None):
    """添加圆角矩形作为内容衬底"""
    from pptx.oxml.ns import qn
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # 调整圆角
    shape.adjustments[0] = 0.02
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=18, bold=False, color=DARK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
             font_name="Microsoft YaHei", level=0):
    """在已有文本框内添加段落"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p

def add_run(p, text, font_size=18, bold=False, color=DARK,
            font_name="Microsoft YaHei"):
    """在段落中添加 run"""
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return run

def add_bottom_bar(slide):
    """添加底部装饰条"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        prs.slide_width, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

def add_top_accent(slide):
    """添加顶部装饰线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def make_slide(title_text, subtitle_text=None):
    """创建一张标准内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_bg(slide, LIGHT)
    add_top_accent(slide)
    add_bottom_bar(slide)

    # 白色内容区域
    content_left = Inches(0.6)
    content_top = Inches(0.5)
    content_width = Inches(12.1)
    content_height = Inches(6.4)
    add_shape_bg(slide, WHITE, content_left, content_top, content_width, content_height)

    # 标题
    add_textbox(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8),
                title_text, font_size=30, bold=True, color=PRIMARY)

    if subtitle_text:
        add_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.5),
                    subtitle_text, font_size=16, color=RGBColor(0x7F, 0x8C, 0x8D))

    return slide

def slide_with_title(title):
    """创建幻灯并返回 (slide, text_frame_for_content)"""
    slide = make_slide(title)
    # 在白色区域内添加内容文本框
    txBox = slide.shapes.add_textbox(
        Inches(1.1), Inches(1.7), Inches(11.1), Inches(5.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return slide, tf

def add_bullet(tf, text, level=0, bold_prefix=None, font_size=17):
    """添加带可选加粗前缀的列表项"""
    if bold_prefix:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        p.level = level
        run = p.add_run()
        run.text = bold_prefix
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = DARK
        run.font.name = "Microsoft YaHei"
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = DARK
        run2.font.name = "Microsoft YaHei"
    else:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        p.level = level
    return p

def add_formula_box(slide, left, top, width, height, formula_text, font_size=16):
    """添加带浅色背景的公式框"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MATH_BG
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)
    shape.adjustments[0] = 0.04

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = formula_text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = "Cambria Math"
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    return shape


# ════════════════════════════════════════════════
# Slide 1: 标题页
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

# 装饰圆（使用半透明颜色）
for x, y, r_val in [
    (Inches(10), Inches(-1), Inches(5)),
    (Inches(-1.5), Inches(4), Inches(3.5)),
    (Inches(11), Inches(5.5), Inches(2)),
]:
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, r_val * 2, r_val * 2
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1E, 0x57, 0x9A)
    circle.line.fill.background()

# 标题文字
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
            "数值最优化算法与理论", font_size=42, bold=True,
            color=WHITE, alignment=PP_ALIGN.LEFT)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.0),
            "第 1 章  引言", font_size=36, bold=True,
            color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.LEFT)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
            "最优化问题概述 · 凸集与凸函数 · 收敛性分析", font_size=18,
            color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.LEFT)

# 底部信息
add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
            "数值最优化算法与理论 — 课程笔记整理", font_size=14,
            color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════
# Slide 2: 目录
# ════════════════════════════════════════════════
slide = make_slide("目  录")
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

items = [
    ("1", "最优化问题概述", "数学模型 · 最优解分类 · 三大规划问题"),
    ("2", "基础知识回顾", "梯度 · Hessian 矩阵 · Taylor 展开 · 向量范数"),
    ("3", "凸集", "定义 · 性质 · 锥与凸锥 · 极点与方向"),
    ("4", "凸函数", "定义 · 等价条件 · 严格凸 · 一致凸"),
    ("5", "凹函数与水平集", "凹函数 · 凸函数的上图像 · 水平集凸性"),
    ("6", "收敛性概念", "局部收敛 · 全局收敛 · 收敛阶"),
    ("7", "附录", "Sherman-Morrison 公式 · 行列式引理"),
]

for i, (num, title, desc) in enumerate(items):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.space_before = Pt(10)
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = f"  ●  {title}"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    run.font.name = "Microsoft YaHei"

    run3 = p.add_run()
    run3.text = f"    —  {desc}"
    run3.font.size = Pt(15)
    run3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
    run3.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 3: 1-1 最优化问题定义
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  最优化问题概述")

p = tf.paragraphs[0]
p.text = "最优化问题的数学模型"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.7),
                "  min  f(x),    x ∈ D ⊆ ℝⁿ        (*)", 20)

add_bullet(tf, "其中 min 是 minimizing (极小化) 的简称", bold_prefix="说明：")
add_bullet(tf, "若 f 为效益函数，则模型为 max f(x) = -min[-f(x)]", level=0)

add_bullet(tf, "f — 目标函数 (Objective Function)")
add_bullet(tf, "D — 可行域 (Feasible Region)")
add_bullet(tf, "x ∈ D — 可行点 (Feasible Point)")

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "通俗理解：在众多可行方案中寻求最佳方案"
run.font.size = Pt(16)
run.font.italic = True
run.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 4: 无约束 vs 约束
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  无约束与约束最优化问题")

p = tf.paragraphs[0]
p.text = "无约束最优化问题 (Unconstrained Optimization)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(5.5), Inches(0.6),
                "  min  f(x),    x ∈ ℝⁿ", 18)

p2 = tf.add_paragraph()
p2.space_before = Pt(16)
run = p2.add_run()
run.text = "约束最优化问题 (Constrained Optimization)"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(4.2), Inches(11), Inches(1.0),
                "  min  f(x)    s.t.   gᵢ(x) ≥ 0,  i ∈ I\n"
                "                     hⱼ(x) = 0,   j ∈ E", 18)

add_bullet(tf, "gᵢ — 不等式约束函数；hⱼ — 等式约束函数")
add_bullet(tf, "可行域 D = { x | gᵢ(x) ≥ 0, hⱼ(x) = 0 }")


# ════════════════════════════════════════════════
# Slide 5: 最优解分类
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  最优解的分类")

p = tf.paragraphs[0]
p.text = ""
# 表格风格展示
definitions = [
    ("局部最优解",
     "设 x* ∈ D，若 ∃ U(x*) 使得 f(x*) ≤ f(x), ∀x ∈ D∩U(x*)"),
    ("严格局部最优解",
     "对 ∀x ∈ D∩U(x*)\\ {x*} 成立严格不等式 f(x*) < f(x)"),
    ("全局(整体)最优解",
     "对 ∀x ∈ D 成立 f(x*) ≤ f(x)"),
    ("严格全局最优解",
     "对 ∀x ∈ D\\ {x*} 成立严格不等式 f(x*) < f(x)"),
]

for i, (term, desc) in enumerate(definitions):
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")
    if i < len(definitions) - 1:
        tf.add_paragraph().space_before = Pt(2)

p2 = tf.add_paragraph()
p2.space_before = Pt(12)
run = p2.add_run()
run.text = "→ 最优化主要研究内容：求解最优解（包括局部与全局）"
run.font.size = Pt(16)
run.font.italic = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 6: 三大规划
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  线性规划 · 二次规划 · 凸规划")

items = [
    ("线性规划 (LP)", "目标函数 f 与约束函数 gᵢ, hⱼ 均为线性函数"),
    ("二次规划 (QP)", "目标函数 f 为二次函数，约束为线性函数"),
    ("凸规划", "f 为凸函数，可行域 D 为凸集"),
]
for term, desc in items:
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "注：凸规划中局部最优解 = 全局最优解"
run.font.size = Pt(16)
run.font.italic = True
run.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 7: 梯度与 Hessian
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  基础知识 — 梯度与 Hessian 矩阵")

p = tf.paragraphs[0]
p.text = "设 f : ℝⁿ → ℝ 二次连续可微"
p.font.size = Pt(18)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.3), Inches(5.5), Inches(1.2),
                "  ∇f(x) = [ ∂f/∂x₁, ∂f/∂x₂, …, ∂f/∂xₙ ]ᵀ", 17)

add_formula_box(slide, Inches(6.8), Inches(2.3), Inches(5.4), Inches(1.2),
                "  ∇²f(x) = [ ∂²f/∂xᵢ∂xⱼ ]   (Hessian 矩阵)", 17)

add_bullet(tf, "一元辅助函数：φ(t) = f[ y + t(x - y) ]")
add_bullet(tf, "一阶导数：φ'(t) = ∇f[ y + t(x-y) ]ᵀ (x-y)")
add_bullet(tf, "二阶导数：φ''(t) = (x-y)ᵀ ∇²f[ y + t(x-y) ] (x-y)")


# ════════════════════════════════════════════════
# Slide 8: Taylor 展开
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  多元函数的 Taylor 展开")

p = tf.paragraphs[0]
p.text = "一阶 Taylor 展开（一阶中值定理）"
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.8),
                "  f(x) = f(y) + ∇f[ y + θ(x-y) ]ᵀ (x-y)  "
                "= f(y) + ∇f(y)ᵀ (x-y) + o(‖x-y‖)", 17)

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "二阶 Taylor 展开（二阶中值定理）"
run.font.size = Pt(19)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(1.2),
                "  f(x) = f(y) + ∇f(y)ᵀ (x-y) + ½(x-y)ᵀ ∇²f[ y+θ(x-y) ] (x-y)\n"
                "       = f(y) + ∇f(y)ᵀ (x-y) + ½(x-y)ᵀ ∇²f(y) (x-y) + o(‖x-y‖²)", 17)

add_bullet(tf, "其中 θ ∈ (0, 1)")


# ════════════════════════════════════════════════
# Slide 9: 向量值函数与范数
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-1  向量值函数的中值定理与范数")

p = tf.paragraphs[0]
p.text = ""
add_bullet(tf, "设 F : ℝⁿ → ℝᵐ 连续可微，F'(x) 为 Jacobi 矩阵", bold_prefix="向量值函数：")

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.8),
                "  F(x) = F(y) + F'(y)(x-y) + o(‖x-y‖)", 18)

add_bullet(tf, "Euclid 范数：‖x‖ = (xᵀx)¹ᐟ²", bold_prefix="向量范数：")
add_bullet(tf, "矩阵从属范数：‖A‖ = max_{‖x‖=1} ‖Ax‖", bold_prefix="矩阵范数：")

add_formula_box(slide, Inches(1.1), Inches(5.0), Inches(11), Inches(0.7),
                "  Frobenius 范数：‖A‖_F = [tr(AAᵀ)]¹ᐟ² = (∑ᵢⱼ aᵢⱼ²)¹ᐟ²", 17)


# ════════════════════════════════════════════════
# Slide 10: 凸集
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  凸集 (Convex Set)")

p = tf.paragraphs[0]
p.text = "凸集的定义"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(5.5), Inches(0.6),
                "  ∀x,y∈S, ∀α∈[0,1]:  αx+(1-α)y ∈ S", 19)

add_bullet(tf, "直观理解：集合包含其中任意两点的连线", bold_prefix="几何意义：")

p2 = tf.add_paragraph()
p2.space_before = Pt(12)
run = p2.add_run()
run.text = "性质"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_bullet(tf, "若 S 是凸集，则 αS 也是凸集（∀α∈ℝ）")
add_bullet(tf, "若 S₁, S₂ 是凸集，则 S₁∩S₂, S₁+S₂, S₁-S₂ 也都是凸集")


# ════════════════════════════════════════════════
# Slide 11: 锥、顶点、方向
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  锥 · 顶点 · 方向")

items = [
    ("锥 (Cone)", "∀λ ≥ 0, ∀x ∈ C: λx ∈ C"),
    ("凸锥 (Convex Cone)", "既是锥又是凸集"),
    ("顶点/极点 (Vertex/Extreme Point)",
     "x ∈ S 不能表示为 S 中两个不同点的凸组合"),
    ("方向 (Direction)", "d ≠ 0 满足 ∀x ∈ S, ∀α ≥ 0: x + αd ∈ S"),
    ("极方向 (Extreme Direction)",
     "d 不能表示为其他两个不同方向的正线性组合"),
]

for term, desc in items:
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")


# ════════════════════════════════════════════════
# Slide 12: 凸函数定义
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  凸函数 (Convex Function)")

p = tf.paragraphs[0]
p.text = ""
add_bullet(tf, "设 S ⊆ ℝⁿ 是凸集，f : ℝⁿ → ℝ", bold_prefix="凸函数定义：")

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.8),
                "  f[αx + (1-α)y] ≤ αf(x) + (1-α)f(y),  ∀x,y∈S, ∀α∈[0,1]", 18)

defs = [
    ("严格凸函数", "∀x≠y, α∈(0,1): 严格不等式 < 成立"),
    ("一致凸/强凸函数",
     "∃m>0: f[αx+(1-α)y] ≤ αf(x)+(1-α)f(y) - mα(1-α)‖x-y‖²"),
]
for term, desc in defs:
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")

p2 = tf.add_paragraph()
p2.space_before = Pt(6)
run = p2.add_run()
run.text = "一致凸 ⇒ 严格凸 ⇒ 凸"
run.font.size = Pt(16)
run.font.italic = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 13: 凸函数的等价条件
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  凸函数的等价条件 (定理 1.2.1)")

p = tf.paragraphs[0]
p.text = "设 f 二次连续可微，则以下命题等价："
p.font.size = Pt(18)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

conditions = [
    "(1) f 是凸函数",
    "(2) φ(t) = f[ t x + (1-t)y ] 是 [0,1] 上的凸函数",
    "(3) f(x) - f(y) ≥ ∇f(y)ᵀ (x-y),  ∀x,y",
    "(4) [∇f(x) - ∇f(y)]ᵀ (x-y) ≥ 0  (梯度单调增)",
    "(5) ∇²f(x) 半正定, ∀x",
]

for cond in conditions:
    add_bullet(tf, cond)

p2 = tf.add_paragraph()
p2.space_before = Pt(6)
run = p2.add_run()
run.text = "→ 该定理是凸函数理论的核心结果之一"
run.font.size = Pt(15)
run.font.italic = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 14: 等价条件证明思路
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  等价条件证明思路")

p = tf.paragraphs[0]
p.text = "(1) ⇔ (2)：利用凸函数定义与一元辅助函数 φ(t) 直接验证"
p.font.size = Pt(17)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

proof_steps = [
    "(1) ⇒ (3)：由凸性得差商不等式，令 α → 0⁺ 取极限",
    "(3) ⇒ (4)：交换 x, y 写出两个不等式，相加即得",
    "(4) ⇒ (5)：令 y = x + tp，利用极限推导得 pᵀ∇²f(x)p ≥ 0",
    "(5) ⇒ (1)：利用二阶 Taylor 展开 + 半正定性得凸性定义",
]
for step in proof_steps:
    add_bullet(tf, step)

p2 = tf.add_paragraph()
p2.space_before = Pt(8)
run = p2.add_run()
run.text = "证明逻辑闭环：(1)→(3)→(4)→(5)→(1)，且(1)⇔(2)"
run.font.size = Pt(16)
run.font.italic = True
run.font.color.rgb = PRIMARY
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 15: 上图像特征
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  凸函数的上图像特征 (定理 1.2.2)")

p = tf.paragraphs[0]
p.text = "定理：设 S 是凸集，则 f 是凸函数 ⟺ f 的上图像是凸集"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.7),
                "  上图像:  P(f) = { (x,α) | x∈S, α∈ℝ, f(x) ≤ α }", 18)

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "必要性证明思路"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"
run2 = p2.add_run()
run2.text = "：若 f 凸，对任意 (x⁽¹⁾,α₁), (x⁽²⁾,α₂) ∈ P(f)，由凸性定义可得凸组合仍在 P(f) 中"
run2.font.size = Pt(16)
run2.font.color.rgb = DARK
run2.font.name = "Microsoft YaHei"

p3 = tf.add_paragraph()
p3.space_before = Pt(6)
run3 = p3.add_run()
run3.text = "充分性证明思路"
run3.font.size = Pt(18)
run3.font.bold = True
run3.font.color.rgb = ACCENT
run3.font.name = "Microsoft YaHei"
run4 = p3.add_run()
run4.text = "：若 P(f) 凸，取 (x⁽¹⁾,f(x⁽¹⁾)), (x⁽²⁾,f(x⁽²⁾)) ∈ P(f)，由凸集定义即得 f 的凸性"
run4.font.size = Pt(16)
run4.font.color.rgb = DARK
run4.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 16: 凹函数与水平集
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  凹函数与水平集")

p = tf.paragraphs[0]
p.text = "凹函数 (Concave Function) — 定理 1.2.3"
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

add_bullet(tf, "定义：-f 是凸函数 ⟺ f 是凹函数")
add_bullet(tf, "等价条件：f(x)-f(y) ≤ ∇f(y)ᵀ(x-y), ∇²f(x) 半负定")

p2 = tf.add_paragraph()
p2.space_before = Pt(12)
run = p2.add_run()
run.text = "水平集的凸性 — 定理 1.2.4"
run.font.size = Pt(19)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(0.6),
                "  水平集:  S_α = { x∈S | f(x) ≤ α }  是凸集", 18)

add_bullet(tf, "证明：对任意 x⁽¹⁾, x⁽²⁾ ∈ S_α，由凸性得凸组合仍在 S_α 中")


# ════════════════════════════════════════════════
# Slide 17: 严格凸函数
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  严格凸函数 (定理 1.2.5)")

p = tf.paragraphs[0]
p.text = "定理：设 f 二次连续可微，下列之一成立则 f 严格凸："
p.font.size = Pt(17)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

add_bullet(tf, "f(x) - f(y) > ∇f(y)ᵀ (x-y),  ∀x≠y", bold_prefix="条件 (1)：")
add_bullet(tf, "∇²f(x) 正定,  ∀x", bold_prefix="条件 (2)：")

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "证明思路"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_bullet(tf, "(2) ⇒ (1)：由二阶 Taylor 展开 + Hessian 正定性得严格不等式")
add_bullet(tf, "(1) ⇒ 严格凸：对任意 x≠y, α∈(0,1)，令 z = αx+(1-α)y，利用 (1) 推导凸组合不等式")

p3 = tf.add_paragraph()
p3.space_before = Pt(6)
run = p3.add_run()
run.text = "注：∇²f(x) 正定是严格凸的充分条件，但不是必要条件"
run.font.size = Pt(15)
run.font.italic = True
run.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 18: 一致凸函数
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  一致凸函数 (定理 1.2.6)")

p = tf.paragraphs[0]
p.text = ""
add_bullet(tf, "f 二次连续可微，则 f 一致凸 ⟺ ∇²f(x) 一致正定", bold_prefix="充要条件：")

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.7),
                "  ∇²f(x) 一致正定  ⟺  ∃m>0: pᵀ∇²f(x)p ≥ m‖p‖², ∀p∈ℝⁿ", 18)

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "必要性证明思路"
run.font.size = Pt(17)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"
run2 = p2.add_run()
run2.text = "：设 f 一致凸，利用一致凸定义 + Taylor 展开 + t→0⁺ 极限"
run2.font.size = Pt(16)
run2.font.color.rgb = DARK
run2.font.name = "Microsoft YaHei"

p3 = tf.add_paragraph()
run3 = p3.add_run()
run3.text = "充分性证明思路"
run3.font.size = Pt(17)
run3.font.bold = True
run3.font.color.rgb = ACCENT
run3.font.name = "Microsoft YaHei"
run4 = p3.add_run()
run4.text = "：设 ∇²f 一致正定，由 Taylor 展开 + 一致正定性得一致凸定义"
run4.font.size = Pt(16)
run4.font.color.rgb = DARK
run4.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 19: 一致凸等价条件
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  一致凸函数的等价条件 (定理 1.2.7)")

p = tf.paragraphs[0]
p.text = "定理：设 f 连续可微，则以下等价："
p.font.size = Pt(18)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

conditions2 = [
    "(1) f 是一致凸函数",
    "(2) ∃m>0: f(x)-f(y) ≥ ∇f(y)ᵀ(x-y) + m‖x-y‖²",
    "(3) ∃m̄>0: [∇f(x)-∇f(y)]ᵀ(x-y) ≥ m̄‖x-y‖²",
]
for cond in conditions2:
    add_bullet(tf, cond)

p2 = tf.add_paragraph()
p2.space_before = Pt(8)
run = p2.add_run()
run.text = "→ 一致凸性保证了强单调梯度，是算法收敛性分析的重要工具"
run.font.size = Pt(15)
run.font.italic = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 20: 收敛性概念
# ════════════════════════════════════════════════
slide, tf = slide_with_title("§1-2  收敛性概念")

p = tf.paragraphs[0]
p.text = "算法收敛性的定义"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = "Microsoft YaHei"

conv = [
    ("局部收敛", "∃U(x*): 初始点 x⁽⁰⁾∈U(x*) 时算法收敛到 x*"),
    ("全局收敛", "对任意初始点 x⁽⁰⁾ 算法收敛到某点 x*"),
]

for term, desc in conv:
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "收敛阶（设 {x⁽ᵏ⁾} → x*）"
run.font.size = Pt(19)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

rates = [
    ("线性收敛", "‖x⁽ᵏ⁺¹⁾-x*‖ ≤ ρ ‖x⁽ᵏ⁾-x*‖,  ρ∈(0,1)"),
    ("超线性收敛", "lim ‖x⁽ᵏ⁺¹⁾-x*‖ / ‖x⁽ᵏ⁾-x*‖ = 0"),
    ("二次收敛", "‖x⁽ᵏ⁺¹⁾-x*‖ ≤ M ‖x⁽ᵏ⁾-x*‖²"),
]
for term, desc in rates:
    add_bullet(tf, desc, bold_prefix=f"▸ {term}：")

p3 = tf.add_paragraph()
run = p3.add_run()
run.text = "二次 ⇒ 超线性 ⇒ 线性"
run.font.size = Pt(15)
run.font.italic = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# Slide 21: Sherman-Morrison
# ════════════════════════════════════════════════
slide, tf = slide_with_title("附录  Sherman-Morrison 公式")

p = tf.paragraphs[0]
p.text = "定理 1.2.8：设 A 非奇异, u,v ∈ ℝⁿ, 1+vᵀA⁻¹u ≠ 0，则"
p.font.size = Pt(17)
p.font.color.rgb = DARK
p.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(2.5), Inches(11), Inches(0.7),
                "  (A + uvᵀ)⁻¹ = A⁻¹ − (A⁻¹uvᵀA⁻¹) / (1 + vᵀA⁻¹u)", 19)

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "引理 1.2.1（行列式计算公式）"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Microsoft YaHei"

add_formula_box(slide, Inches(1.1), Inches(4.2), Inches(11), Inches(0.7),
                "  det(I+u₁v₁ᵀ) = 1+u₁ᵀv₁", 18)
add_formula_box(slide, Inches(1.1), Inches(5.1), Inches(11), Inches(0.7),
                "  det(I+u₁v₁ᵀ+u₂v₂ᵀ) = (1+u₁ᵀv₁)(1+u₂ᵀv₂) − (u₁ᵀv₂)(v₁ᵀu₂)", 18)


# ════════════════════════════════════════════════
# Slide 22: 总结
# ════════════════════════════════════════════════
slide, tf = slide_with_title("本章总结")

p = tf.paragraphs[0]
p.text = ""
summary_items = [
    "建立了最优化问题的数学模型，区分了无约束/约束问题",
    "定义并分类了局部最优解、全局最优解等概念",
    "回顾了梯度、Hessian 矩阵、Taylor 展开等分析工具",
    "系统介绍了凸集的定义、性质及锥、顶点、方向等概念",
    "深入讨论了凸函数、严格凸函数、一致凸函数的定义与等价条件",
    "介绍了算法的收敛性概念（局部/全局收敛、线性/超线性/二次收敛）",
    "给出 Sherman-Morrison 公式及其相关行列式引理",
]

for item in summary_items:
    add_bullet(tf, item)

p2 = tf.add_paragraph()
p2.space_before = Pt(10)
run = p2.add_run()
run.text = "→ 本章为后续各章算法分析奠定了理论基础"
run.font.size = Pt(17)
run.font.bold = True
run.font.color.rgb = PRIMARY
run.font.name = "Microsoft YaHei"


# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
output_path = r"C:\Users\MoMoe\Desktop\project\documents\第1章_引言.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
